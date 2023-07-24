#提取所有文本字符
from pptx import Presentation
data = []
prs = Presentation('data\制造业必修课.pptx')
for slide in prs.slides: #遍历每页PPT
    for shape in slide.shapes: #遍历PPT中的每个形状
        if shape.has_text_frame: #判断该是否包含文本，保证有文本才提取
            for paragraph in shape.text_frame.paragraphs: #按文本框中的段落提取
                data.append(paragraph.text) #提取一个段落的文本，就存到列表data中


#写入文本文件
TxtFile = open('data\制造业必修课.txt', 'w',encoding='utf-8')
for i in data:
    TxtFile.write(i+'\n') #写入并换行，以保证正确分段
TxtFile.close() #保存


#写入word文件
import docx
doc=docx.Document()#创建一个word文件对象
for i in data:
    doc.add_paragraph(i) #增加一个段落，并将列表中的一个字符串写入word文件
doc.save('data\制造业必修课.docx')#保存
