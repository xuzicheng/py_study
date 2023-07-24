from pptx import Presentation
prs = Presentation('data\荣誉证书模板.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[0]) # 第一个模板的第0个板式
for ph in slide.placeholders: #遍历这页PPT的所有占位符
    phf = ph.placeholder_format #获取占位符的格式
    print(phf.idx) #打印其ID编号
    ph.text = str(phf.idx)# 将编号写入PPT对应的位置中，以便后面一一对应
# 以上读取到占位符的ID方便后面调用
prs.save('data\荣誉证书模板-占位符编号.pptx')

from openpyxl import load_workbook
wb = load_workbook("data\名单.xlsx")
ws = wb.active
data={}
for row in range(2,ws.max_row+1):
    class_id = ws['A' + str(row)].value
    name = ws['B' + str(row)].value 
    data.setdefault(class_id,[])
    data[class_id].append(name)


import time
t0=time.time()# 程序开始运行的时间

prs = Presentation('data\荣誉证书模板.pptx')
slide_layout = prs.slide_layouts[0] #调用设置好的母版，因为是母版的第一版式，所以取[0]
for class_id in data:
    for name in data[class_id]:
        slide = prs.slides.add_slide(slide_layout) #以母版的版式为基础新增一页幻灯片
        #往幻灯片中写入内容
        slide.placeholders[10].text = "{} 班 {} 同学：".format(class_id,name) #此处是班级和姓名
        slide.placeholders[11].text = "在2019-2020学年度第一学期获得"
        slide.placeholders[12].text = "“好孩子”称号。"
        slide.placeholders[13].text = "特发此证，以资鼓励。"
        slide.placeholders[14].text = "市幼儿园"
        slide.placeholders[15].text = "2020年1月"
prs.save('data\荣誉证书(总).pptx')
t1 = time.time()

print('程序用时：',str(round(t1-t0))+'秒。')